import pandasai as pai #for future use: use dataframe instead of smart dataframe: smart dataframe will be depricated soon
from pandasai_litellm.litellm import LiteLLM
from pandasai import SmartDataframe
from pandasai.config import Config
import pandas as pd
from canvasapi import Canvas #for canvas data
from canvasapi.folder import Folder
from pptx import Presentation #parse and covert powerpoints to strings
from docx import Document #parse and convert docx files to strings
from PyPDF2 import PdfReader #convert pdf to strings
from bs4 import BeautifulSoup #For Syllabus body
from prompt_toolkit.shortcuts import ProgressBar #For loading Bars
import os
import io
import re
def assignments(courses): #Pulls assignments from canvasapi could be in main but I like it clean
    course_assignments = []
    all_courses = []
    with ProgressBar() as progressor:
        for course in progressor(courses, label="Loading Assignments..."):
            course_assignments = course.get_assignments(include=['rubric', 'rubric_settings', 'submission']) #Get needed data with each assignment
            all_courses.append(course_assignments)
  
    ''' For debugging
    for course in all_courses:
        for i in course:
            print(i.name)
    '''
    return all_courses
def assign_dicts(course ,assignment): #stores a assignement in a dictionary
    #Default Values
    graded = False
    score = None
    rubric_id = "None"
    rubric = "No Rubric"
    if hasattr(assignment, 'submission'): #Check if assignment is graded
        sub = assignment.submission
        if sub.get('workflow_state') == 'graded':
            score = sub.get('score')
            if score:
                score = float(score)
            graded = True
    if hasattr(assignment, 'rubric_settings') or hasattr(assignment, 'rubric'): #Try to get rubric not all assignments have one
            rubric_id = assignment.rubric_settings.get('id')
            if rubric_id:
                rubric_list = []
                for fields in assignment.rubric: #rubric is a list of dictionaries
                    desc = fields.get('description', 'No Description')
                    points = fields.get('points', 0)
                    rubric_list.append(f".Desc: {desc}, points: {points}")
                rubric = "\n".join(rubric_list)
                #print(f"{assignment.name} {rubric}")
    
    max_points = getattr(assignment, "points_possible", None) #Get possible points
    if max_points:
        max_points = float(max_points)
        
    course_name = course.name
    dict_assignment = {
        'course_name': course_name,
        'id': str(getattr(assignment, "id", None)),
        'name': getattr(assignment, "name", None),
        'url': getattr(assignment, "html_url", None),
        'prob_description': getattr(assignment, "description", None),
        'max_points': max_points,
        'due_date': getattr(assignment, "due_at", None),
        'late_end_date': getattr(assignment, "lock_at", None),
        'type': getattr(assignment, "grading_type", None),
        'sub_type': ", ".join(getattr(assignment, "grading_type", None)),
        'timezone': "UTC",
        'score': score,
        'grade_state': graded,
        'rubric_exists': str(rubric_id),
        'rubric': rubric
    }
    return dict_assignment
def assignment_dict_list(all_courses, assignment_list): #Assembles a list of Assignment dictionaries
                                        #Could easily be in main but I like keeping it clean
    assignment_dictions = []
    i = 0
    with ProgressBar() as progressor:
        for course in progressor(assignment_list, label="Storing File Data..."):
            for assingment in progressor(list(course), label=f"Loading {all_courses[i].name} Files..."):
                assignment_dictions.append(assign_dicts(all_courses[i], assingment))
            i += 1
    return assignment_dictions
def bytes_text(bytes, filename):
    file_contents = []
    content_text = ""
    if filename.endswith(".pdf"):
        file = PdfReader(io.BytesIO(bytes))
        for content in file.pages:
            content_text += str(content.extract_text)
    elif filename.endswith(".pptx"):
        pres = Presentation(io.BytesIO(bytes))
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    content_text += shape.text.strip()
    elif filename.endswith((".docx", ".doc")):
        file = Document(io.BytesIO(bytes))
        for content in file.paragraphs:
            content_text += content.text
    else:
        try:
            content_text = bytes.decode("utf-8", errors="ignore")
        except Exception:
            return [] 
    file_list = re.split(r"\.|\?|!|\n", content_text)
    content_number = 1
    for content_text in file_list:
        file_contents.append({'content_number': content_number, 'content_text': content_text})
        content_number += 1
    return file_contents
def get_files(all_courses): #Extracts files from courses
    #Numerous try except blocks for cases where the user isn't authorized to retrieve data
    #from canvasapi (very common for a student user in my experience)
    course_names = []
    course_syllabi = []
    file_contents = []
    course_files = []
    i = 0
    with ProgressBar() as progressor:   
        for course in progressor(list(all_courses), label="Reading Files"):
            course_names.append(course.name)
            fol_number = 1
            syllabus = None
            try:
                folders = list(course.get_folders())
            except Exception:
                #debugging print(e)
                continue
            for folder in progressor(list(folders), label=f"{all_courses[i].name} Folders..."):
                if getattr(folder, 'hidden', False):
                    continue
                try:
                    files = list(folder.get_files())
                except Exception:
                    #Debugging print(e)
                    continue
                for file in progressor(list(files), label=f"Folder: {fol_number}..."):
                    if getattr(file, 'hidden_for_user', False) or getattr(file, 'locked', False) or isinstance(file, Folder):
                        continue    
                    file_contents = file_dicts(course, file)
                    filename = str(file.display_name).lower()
                    if "syllabus" in filename:
                        syllabus = {'filename': filename, 'url': file.url}  
                    if file_contents != []:
                        course_files.append(file_contents)
                    else:
                        continue
                fol_number += 1
            i += 1
            course_syllabi.append(syllabus)
    return course_names, course_syllabi, course_files
def files_rows(course_files): #Prepares files for dataframe 
    file_rows = []
    with ProgressBar() as progressor:
        for file in progressor(course_files, label="Converting files to Dataframe..."):
            file_text = file.get('file_text')
            if isinstance(file_text, list): #Case for a list of file dictionaries
                for content in file_text:
                    data_text = content.get('content_text')
                    if data_text:
                        file_rows.append({
                            "course_name": file.get('course_name'),
                            "filename": file.get('filename'),
                            "text_sentance_index": content.get('slide_number'),
                            "file_sentance_text": data_text,
                            #"full_text": content.get('full_text')
                        })
            else: #Case for a previously invalid file in dictionize files
                file_rows.append({
                    "course_name": file.get("course_name"),
                    "filename": file.get("filename"),
                    "text_sentance_index": None,
                    "file_sentance_text": None,
                    #"full_text": None
                })
    return file_rows
def file_dicts(course, course_file): #Put file information into a dictionary
    #Default values
    contents = "Download failed"
    filename = "Unaccessable"
    type = "None"
    content_list = []
    url = "None"
   
    contents = course_file.get_contents(binary=True) #Get binary for file
    type = getattr(course_file, 'content-type', None)
    url = getattr(course_file,'url', None)
    
    filename = course_file.display_name
    if filename: 
        str_filename = str(filename).lower()
        file_contents = bytes_text(contents, str_filename) #Send off file data to convert to text
        if file_contents:
            for content in file_contents:
                content_list.append({'content_number': content['content_number'], 'content_text': content['content_text']})
        else: #Case for invalid file
                content_list = None
        file_dic = {
            'file_course_name': getattr(course, "name", None),
            'filename': getattr(course_file, "display_name", None),
            'url': url,
            'file_type': type,
            'file_text': content_list,
        }
        return file_dic
    else: #Case for bad file
        return [] 
def syl_dicts(course_syllabi, course_objs): #Converts Syllabi to dictionary for ai
    c_num = 0
    course_rows = []
    with ProgressBar() as progressor:
        for syllabus in progressor(course_syllabi, label="Converting syllabi to Dataframe..."): #unfortunatly can't convert to list to track progress
            course = course_objs[c_num]
            total_syll = getattr(course, 'syllabus_body', None) #Get Syllabus body from canvas
            #Default values
            syllabus_url = None
            syllabus_name = None 
            syllabus_file = False

            if syllabus: #check if syllabus file exists for course and store data
                syllabus_name = syllabus.get('filename')
                syllabus_url = syllabus.get('url')
                syllabus_file = True
            if total_syll: #case for a syllabus page
                html = BeautifulSoup(total_syll, features="lxml") #Need to convert from html since Canvasapi only return the syllabus body as html
                sentences = str(html.get_text())
                sentences = re.split(r"\.|\?|!|\n", sentences) #Divide Syllabus into a list of sentences
                syllabus_sentance = 1
                for sentance in sentences: #store each sentance from Syllabus into a dictionary so ai may summarize it
                    course_rows.append({         
                        'has_syllabus': True,           
                        'syllabus_course_name': course.name,
                        'syllabus_file': syllabus_file,
                        'syllabus_filename': syllabus_name,
                        'syllabus_file_url': syllabus_url,
                        'syllabus_webpage': True,
                        'syllabus_webpage_sentance': sentance.strip(),
                        'syllabus_webpage_index': syllabus_sentance
                        })
                    syllabus_sentance += 1
            else: #case for no syllabus page. Usually a file instead
                course_rows.append({         
                        'has_syllabus': True,           
                        'syllabus_course_name': course.name,
                        'syllabus_file': syllabus_file,
                        'syllabus_filename': syllabus_name,
                        'syllabus_file_url': syllabus_url,
                        'syllabus_webpage': False,
                        'syllabus_webpage_sentance': None,
                        'syllabus_webpage_index': None
                        })
            c_num += 1
    return course_rows
def allowed_courses(courses): #Trims Courses user isn't allowed to access
    all_courses = []
    exceptions = 0
    legal_courses = 0
    courses_list = list(courses) #convert to list so progressbar can track progress
    courses_list = courses
    with ProgressBar() as progressor:
        for course in progressor(courses_list, label="Loading Legal Courses..."):
            try:
                name = course.name
                assert(course.get_assignments())
                #print(f"Name: {name}\n")
                if name != "no-name":
                    all_courses.append(course)
                    legal_courses += 1

            except Exception: #Case for course user can't acess for unt only current courses may be accessed
                exceptions += 1 
                #print("exception\n")
    print(f"Illegal Courses: {exceptions}, Legal Courses: {legal_courses}") #Kinda cool to see how many things you don't have access to
    return all_courses 
def enrollments(all_enrollments): #Function trims enrollments the user isn't authorized to acess in canvas
    valid_enrollments = []   
    for enrollment in all_enrollments:
        try: #needs try except block incase user isn't authorized to get data from canvas
            id = enrollment.course_id
            if id is not None:
                valid_enrollments.append(enrollment)
        except Exception: #case where user isn't authorized to access an enrollment
            pass
    return valid_enrollments #return only valid enrollments
def gather_averages(user_canvas, valid_enrollments): #Get course scores and grades from all enrollments
    course_avrs = []
    i = 0
    for enrollment in valid_enrollments: #loop through enrollments and store score information in a dictionary
        id = enrollment.course_id
        course = user_canvas.get_course(id)
        course_avr = {
            #in my experience UNT only has the final and current scores no letter grade
            'averages_course_name': course.name,
            'current_score_average': enrollment.grades.get('current_score') or None,
            'final_score_average': enrollment.grades.get('final_score') or None,
            'current_grade_letter': enrollment.grades.get('current_grade') or None,
            'final_grade_letter': enrollment.grades.get('final_grade') or None
        }
        course_avrs.append(course_avr)
        i += 1
    return course_avrs #return score information for dataframe

if __name__ == "__main__":
    #Set up canvas object
    CANVAS_URL = "https://unt.instructure.com/"
    CANVAS_API = os.environ.get("CANVAS_API_KEY")
    GEMINI_KEY = os.environ.get("GEMINI_API_KEY")
    if not CANVAS_API and GEMINI_KEY:
        SystemExit("ERROR: Ensure Canvas and GEMINI API keys are stored in system variables as CANVAS_API_KEY and GEMINI_API_KEY")

    try:
        user_canvas = Canvas(CANVAS_URL, CANVAS_API)
        user = user_canvas.get_current_user()
    except Exception:
        print("Canvas authorization failed:")
        raise SystemExit("Check CANVAS_API_KEY yenviromental variable.")
    else:

        courses = user.get_courses(enrollment_status='active', include=['syllabus_body']) #get courses from Canvasapi
        all_courses = allowed_courses(courses) #call allowed_courses to remove unauthorized courses
       #Get files was originally in main which is why there is a tupple I think this looks cleaner in my opinion
        course_names, course_syllabi, course_files = get_files(all_courses) #Get information for course names, copys of syllabi files for course rows and all files
        course_rows = syl_dicts(course_syllabi, all_courses) #Convet Syllabi into dictionaries
        assignment_list = assignments(all_courses) #Get all assignments from all_courses
        assignment_dictions = assignment_dict_list(all_courses, assignment_list) #convert assignments into dictionaries
        file_rows = files_rows(course_files) #prepare files into one list for dataframe
        all_enrollments = user.get_enrollments(include=['total_scores', 'grades'], state=["active"]) #Get enrollments for scores
        allowed_enrollments = enrollments(all_enrollments) #Trim unauthorized enrollments
        average_rows = gather_averages(user_canvas, allowed_enrollments) #Create dictionaires for course grades
    #os.system('cls' if os.name == 'nt' else 'clear')
    INSTRUCTIONS = """
Use only the information that exists in the dataframe: course names, files, assignments, due dates, scores, 
rubrics, syllabus text, and slide/paragraph text materials. Don't guess or make up information.
Keep your answers short and focused. If a users question isn't clear, ask one clarifying question. 
If the needed information truly isn't anywhere in the dataframe, respond with: "No relevant data found."

Quiz, Flashcards, Study Questions:
If the user asks for anything like a quiz, flashcards, study questions, a review sheet,
practice test, or multiple-choice questions, you are allowed to generate new questions—
but every question and answer must be based strictly on text found in the data 
especially the “text” column from files.

Do not make up facts or add outside knowledge.

If the user mentions a specific filename, course, slide, or topic, limit the questions to
only that material.

When creating quizzes:
Aim for 5-20 questions unless the user asks for a different amount. Multiple-choice questions 
should have four options (A-D) with one correct answer. Include an answer key at the end.
Keep the questions clear and related to the text.

These rules must be followed.
"""
    userName = input("Please enter your name: ")
    user_data = {
        'users_name': userName,
        'course_names': course_names,
        'quit_help': "Type Quit to exit program",
    }   

    llm = LiteLLM(model="gemini/gemini-2.5-flash", provider="google_ai_studio", api_key=GEMINI_KEY)
    config = Config(
        llm=llm,
        temperature=0.2,
        seed=26
    )


    #Create dataframes for data collected from canvasapi
    assignment_data = pd.DataFrame(assignment_dictions)
    user_info = pd.DataFrame(user_data)
    files_data = pd.DataFrame(file_rows)
    syllabus_data = pd.DataFrame(course_rows)
    average_data = pd.DataFrame(average_rows)
    #Combine dataframes into one dataframe
    ai_frame = pd.concat([user_info, average_data, syllabus_data, files_data, assignment_data], ignore_index=True)

    dataframe = SmartDataframe(ai_frame, config=config)
    os.system('cls' if os.name == 'nt' else 'clear')
    userinput = input("How may I help you: ")
    while userinput != "quit" and userinput != "Quit":
        try:
            instructs = INSTRUCTIONS + "\nPrompt: " + userinput
            prompt = dataframe.chat(instructs)
            
            print(prompt)
        except Exception as e:
            print(f"AI model currently unavailable: {e}")
        userinput = input("How may I help you: ")
